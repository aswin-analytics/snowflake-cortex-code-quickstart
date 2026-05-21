"""
Generate a visual PowerPoint presentation for the Snowflake Cortex Code Quickstart workshop.

Design:
- Dark theme throughout (consistent navy background)
- Native PPT shapes for flow diagrams (no pixelated images)
- 3-column stackable flows showing the journey on each slide
- Progress bar at the bottom linking slides together
- Meaningful descriptions (not overly minimal)

Usage:
    pip install python-pptx
    python create_presentation.py

Output: Cortex_Code_Demo_Presentation.pptx (in the same directory)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ──
DARK = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x29, 0xB5, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xAA, 0xBB, 0xCC)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
PURPLE = RGBColor(0x9B, 0x59, 0xB6)
RED_ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
GRAY_BOX = RGBColor(0x2C, 0x3E, 0x50)

# Progress bar section labels
SECTIONS = [
    "Intro", "CoCo", "Overview", "Lineage",
    "Demo 1", "Demo 2", "Demo 3",
    "Features", "Cost", "Controls", "Summary"
]


def set_dark_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK


def add_progress_bar(slide, current_index):
    """Add a progress bar at the bottom showing current position."""
    bar_top = Inches(7.0)
    total_width = 12.333
    segment_width = total_width / len(SECTIONS)
    start_left = 0.5

    for i, label in enumerate(SECTIONS):
        left = start_left + i * segment_width
        # Segment rectangle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), bar_top,
            Inches(segment_width - 0.05), Inches(0.35)
        )
        if i == current_index:
            shape.fill.solid()
            shape.fill.fore_color.rgb = BLUE
        elif i < current_index:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0x1E, 0x6F, 0x8F)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = GRAY_BOX
        shape.line.fill.background()

        # Label
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Pt(2)
        tf.margin_right = Pt(2)
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(8)
        p.font.bold = (i == current_index)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER


def add_title(slide, text):
    """Add slide title."""
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT


def add_flow_box(slide, text, sub_text, left, top, width, height, color):
    """Add a rounded rectangle flow box with title and description."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if sub_text:
        p2 = tf.add_paragraph()
        p2.text = sub_text
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        p2.alignment = PP_ALIGN.CENTER


def add_arrow(slide, left, top, width):
    """Add a horizontal arrow connector."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, 
        Inches(left), Inches(top), Inches(width), Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()


def add_down_arrow(slide, left, top, height=0.4):
    """Add a vertical down arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(left), Inches(top), Inches(0.3), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()


def add_text_block(slide, text, left, top, width, size=Pt(15), color=WHITE):
    """Add a text block."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = color
        p.space_after = Pt(8)


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)

txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.3), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Cortex Code for AP Data Engineering"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.3), Inches(1))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "From Raw ERP Data to Natural-Language Analytics"
p2.font.size = Pt(22)
p2.font.color.rgb = BLUE
p2.alignment = PP_ALIGN.CENTER

txBox3 = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(11.3), Inches(1))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "An AI-Powered Workflow using Snowflake Cortex Code"
p3.font.size = Pt(16)
p3.font.color.rgb = LIGHT
p3.alignment = PP_ALIGN.CENTER

# 3-column preview boxes
add_flow_box(slide, "Pipeline Builder", "Dynamic Tables", 1.5, 5.5, 3.0, 0.9, ORANGE)
add_flow_box(slide, "Pipeline Maintenance", "PRD-Driven Evolution", 5.2, 5.5, 3.0, 0.9, BLUE)
add_flow_box(slide, "Cortex Agent", "NL Analytics", 8.9, 5.5, 3.0, 0.9, GREEN)
add_arrow(slide, 4.5, 5.8, 0.7)
add_arrow(slide, 8.2, 5.8, 0.7)

add_progress_bar(slide, 0)

# ============================================================
# SLIDE 2: What is Cortex Code?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)
add_title(slide, "What is Cortex Code?")

# 3-column flow: Entry → Agent → Snowflake
add_flow_box(slide, "Entry Points", "CLI (Terminal/IDE)\nSnowsight (Browser)", 0.6, 1.5, 3.5, 1.8, GRAY_BOX)
add_arrow(slide, 4.1, 2.2, 0.8)
add_flow_box(slide, "CoCo AI Agent", "Plans, writes, and\nexecutes on your behalf", 4.9, 1.5, 3.5, 1.8, GREEN)
add_arrow(slide, 8.4, 2.2, 0.8)
add_flow_box(slide, "Snowflake", "Your data, compute,\nand identity (RBAC)", 9.2, 1.5, 3.5, 1.8, PURPLE)

# Description below
add_text_block(slide,
    "Snowflake's AI coding agent purpose-built for data engineering, analytics, and ML\n"
    "Operates directly inside your Snowflake environment with full role-based access control\n"
    "Understands your data catalog, schema metadata, and Snowflake best practices\n"
    "Executes SQL as your identity — no separate service accounts needed\n"
    "Extensible with reusable Skills, custom Cortex Agents, and MCP integrations",
    0.6, 3.8, 12.0, Pt(15), WHITE)

add_progress_bar(slide, 1)

# ============================================================
# SLIDE 3: Workshop Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)
add_title(slide, "Workshop: AP Invoices Pipeline")

# 3-column flow for the 3 demos
add_flow_box(slide, "Demo 1", "Pipeline Builder\nCreate Dynamic Table\nfrom SAP + Oracle", 0.6, 1.5, 3.5, 2.2, ORANGE)
add_arrow(slide, 4.1, 2.4, 0.8)
add_flow_box(slide, "Demo 2", "Pipeline Maintenance\nOnboard Baan + Workday\nfrom PRD requirements", 4.9, 1.5, 3.5, 2.2, BLUE)
add_arrow(slide, 8.4, 2.4, 0.8)
add_flow_box(slide, "Demo 3", "Cortex Agent\nSemantic View + Agent\nfor NL analytics", 9.2, 1.5, 3.5, 2.2, GREEN)

# Context below
add_text_block(slide,
    "Scenario: Finance needs a trusted, unified AP invoices layer across 4 ERP systems\n"
    "Data sources: SAP (15 invoices), Oracle (15), Baan IV (10), Workday (10)\n"
    "Goal: Standardize, validate, evolve, and make queryable in plain English\n"
    "Single storyline from data discovery through to operationalized AI analytics",
    0.6, 4.2, 12.0, Pt(15), WHITE)

add_progress_bar(slide, 2)

# ============================================================
# SLIDE 4: Pipeline Lineage
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)
add_title(slide, "Pipeline Lineage: End-to-End Data Flow")

# Row 1: Bronze sources (4 boxes)
box_w = 2.8
row1_top = 1.4
add_flow_box(slide, "SAP", "15 invoices", 0.4, row1_top, box_w, 0.9, ORANGE)
add_flow_box(slide, "Oracle", "15 invoices", 3.4, row1_top, box_w, 0.9, ORANGE)
add_flow_box(slide, "Baan IV", "10 invoices", 6.4, row1_top, box_w, 0.9, ORANGE)
add_flow_box(slide, "Workday", "10 invoices", 9.4, row1_top, box_w, 0.9, ORANGE)

# Down arrows
for x in [1.6, 4.6, 7.6, 10.6]:
    add_down_arrow(slide, x, 2.35, 0.35)

# Row 2: Silver layer
add_flow_box(slide, "SILVER_AP_INVOICES", "Dynamic Table — unified schema, ~50 rows\nNormalized status, dedup, column aliasing", 1.5, 2.8, 10.3, 1.2, BLUE)

# Down arrow center
add_down_arrow(slide, 6.5, 4.05, 0.35)

# Row 3: Consumption layer (3 boxes)
row3_top = 4.5
add_flow_box(slide, "Semantic View", "SV_AP_ANALYTICS\n7 metrics, 7 VQRs", 0.6, row3_top, 3.5, 1.2, PURPLE)
add_flow_box(slide, "Cortex Agent", "AP_ANALYTICS_ASSISTANT\nNL-to-SQL via Analyst", 4.9, row3_top, 3.5, 1.2, GREEN)
add_flow_box(slide, "Streamlit App", "Interactive dashboard\n+ AI chat interface", 9.2, row3_top, 3.5, 1.2, GRAY_BOX)

# Arrows between consumption layer
add_arrow(slide, 4.1, 5.0, 0.8)
add_arrow(slide, 8.4, 5.0, 0.8)

# Legend
add_text_block(slide,
    "Bronze (Source Tables)  →  Silver (Dynamic Table)  →  Semantic + AI (Consumption Layer)",
    0.6, 6.2, 12.0, Pt(13), LIGHT)

add_progress_bar(slide, 3)

# ============================================================
# SLIDE 5: Demo 1 — Pipeline Builder
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)
add_title(slide, "Demo 1: Pipeline Builder")

# 3-column flow
add_flow_box(slide, "Data Discovery", "CoCo profiles SAP &\nOracle bronze tables", 0.6, 1.5, 3.5, 1.6, GRAY_BOX)
add_arrow(slide, 4.1, 2.1, 0.8)
add_flow_box(slide, "Plan Mode", "Review DT design\nbefore executing", 4.9, 1.5, 3.5, 1.6, PURPLE)
add_arrow(slide, 8.4, 2.1, 0.8)
add_flow_box(slide, "Dynamic Table", "SILVER_AP_INVOICES\nSAP + Oracle unified", 9.2, 1.5, 3.5, 1.6, BLUE)

# Detailed steps
add_text_block(slide,
    "Challenge: Finance has SAP & Oracle AP data but no standardized Silver layer\n\n"
    "Step 1 — Data Discovery: CoCo explores the schema, profiles columns, identifies mismatches\n"
    "Step 2 — Schema Comparison: Naming differences (INV_ID vs INVOICE_ID), type alignment\n"
    "Step 3 — Plan Mode (Ctrl+P): Full DT design reviewed before any execution\n"
    "Step 4 — Dynamic Table Created: Column aliasing, payment terms normalization (N30→NET30),\n"
    "          status normalization (APPROVED/PENDING/OTHER), SOURCE_SYSTEM tag per branch\n"
    "Step 5 — $dynamic-tables Skill: Validates TARGET_LAG, refresh mode, and best practices",
    0.6, 3.5, 12.0, Pt(14), WHITE)

add_progress_bar(slide, 4)

# ============================================================
# SLIDE 6: Demo 2 — Pipeline Maintenance
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)
add_title(slide, "Demo 2: Pipeline Maintenance — PRD-Driven Evolution")

# 3-column flow (5 steps in a row)
add_flow_box(slide, "PRD Files", "3 CSVs with\nrequirements", 0.3, 1.5, 2.3, 1.4, ORANGE)
add_arrow(slide, 2.6, 2.0, 0.5)
add_flow_box(slide, "Gap Analysis", "Auto-compare\nvs current DT", 3.1, 1.5, 2.3, 1.4, PURPLE)
add_arrow(slide, 5.4, 2.0, 0.5)
add_flow_box(slide, "Open Questions", "Surface ambiguity\nfor human decision", 5.9, 1.5, 2.3, 1.4, RED_ACCENT)
add_arrow(slide, 8.2, 2.0, 0.5)
add_flow_box(slide, "DDL Generated", "4-source DT\nwith dedup logic", 8.7, 1.5, 2.3, 1.4, BLUE)
add_arrow(slide, 11.0, 2.0, 0.5)
add_flow_box(slide, "Validation", "5 proof queries\nconfirm correctness", 11.5, 1.5, 1.5, 1.4, GREEN)

# Detailed steps
add_text_block(slide,
    "Challenge: Onboard 2 new ERP systems (Baan IV + Workday) from business requirements\n\n"
    "PRD Ingestion: CoCo reads 3 CSV files — source onboarding, column mappings, business rules\n"
    "Gap Analysis: Automatically identifies new UNION ALL branches needed, rule changes required\n"
    "Open Questions Surfaced: Payment terms consistency, Baan dedup scope, Workday case sensitivity\n"
    "  — CoCo never guesses — always surfaces ambiguity for human decision before proceeding\n"
    "DDL Generated: Full 4-source Dynamic Table with QUALIFY dedup on Baan, extended status CASE\n"
    "Validation: Record counts, NULL checks, dedup verification, status normalization, payment terms\n"
    "Reusable Skill: $prd-to-dt-plan created — any future PRD gets the same rigorous workflow",
    0.3, 3.3, 12.5, Pt(13), WHITE)

add_progress_bar(slide, 5)

# ============================================================
# SLIDE 7: Demo 3 — Cortex Agent
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)
add_title(slide, "Demo 3: Cortex Agent with Semantic View")

# 3-column flow
add_flow_box(slide, "Semantic View", "Define business meaning\nof every column & metric", 0.6, 1.5, 3.5, 1.6, PURPLE)
add_arrow(slide, 4.1, 2.1, 0.8)
add_flow_box(slide, "Cortex Agent", "NL-to-SQL translation\nvia Cortex Analyst", 4.9, 1.5, 3.5, 1.6, GREEN)
add_arrow(slide, 8.4, 2.1, 0.8)
add_flow_box(slide, "Business Users", "Ask questions in English\nGet answers + charts", 9.2, 1.5, 3.5, 1.6, GRAY_BOX)

# Detailed description
add_text_block(slide,
    "Challenge: Make the Silver layer queryable in plain English for non-SQL business users\n\n"
    "Semantic View (SV_AP_ANALYTICS):\n"
    "  • 11 dimensions (vendor, currency, cost center, GL, status, source system)\n"
    "  • 2 time dimensions (invoice date, due date) for trending and overdue analysis\n"
    "  • 7 pre-built metrics (total spend, invoice count, avg amount, overdue count, pending amount)\n"
    "  • 7 verified queries as ground-truth examples for the AI\n\n"
    "Cortex Agent (AP_ANALYTICS_ASSISTANT):\n"
    "  • Users ask: \"What is total spend by vendor?\" or \"Show overdue invoices\"\n"
    "  • Warns about currency mixing, asks clarifying questions when ambiguous\n"
    "  • Full Streamlit app with KPI cards, charts, and AI chat interface",
    0.6, 3.5, 12.0, Pt(13), WHITE)

add_progress_bar(slide, 6)

# ============================================================
# SLIDE 8: Key Features Demonstrated
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)
add_title(slide, "Key Cortex Code Features Demonstrated")

# Left column: Capabilities
txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(5.5), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Core Capabilities"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = BLUE

add_text_block(slide,
    "Data Discovery — #table references, DESCRIBE, profiling\n"
    "Plan Mode (Ctrl+P) — review and approve before executing\n"
    "Bundled Skills — $dynamic-tables for DT analysis & best practices\n"
    "File Ingestion — @file to load CSVs/XLSX into conversation context\n"
    "Custom Skills — $prd-to-dt-plan for repeatable PRD workflows\n"
    "Semantic View Authoring — YAML-based business model definition\n"
    "Cortex Agent Deployment — NL-to-SQL analytics assistant",
    0.6, 1.9, 5.8, Pt(14), WHITE)

# Right column: Business Value
txBox2 = slide.shapes.add_textbox(Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.5))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Business Value"
p2.font.size = Pt(18)
p2.font.bold = True
p2.font.color.rgb = BLUE

add_text_block(slide,
    "Accelerates pipeline development (hours, not days)\n"
    "Enforces consistency via repeatable skills across team\n"
    "Surfaces ambiguity — never makes silent assumptions\n"
    "Full audit trail — plans, handoff docs, proof queries\n"
    "Democratizes data access for non-SQL business users\n"
    "Reduces risk with Plan Mode + validation steps\n"
    "Reusable artifacts for future projects and teams",
    7.0, 1.9, 5.8, Pt(14), WHITE)

add_progress_bar(slide, 7)

# ============================================================
# SLIDE 9: Cost Model
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)
add_title(slide, "Cortex Code Cost Model")

# 3-column flow: Tokens → Credits → Billing
add_flow_box(slide, "Your Prompt", "Input tokens sent\nto the AI model", 0.6, 1.5, 3.5, 1.5, BLUE)
add_arrow(slide, 4.1, 2.0, 0.8)
add_flow_box(slide, "AI Inference", "Model processes and\ngenerates response", 4.9, 1.5, 3.5, 1.5, GREEN)
add_arrow(slide, 8.4, 2.0, 0.8)
add_flow_box(slide, "Credits Charged", "Per million tokens\n(input + output)", 9.2, 1.5, 3.5, 1.5, ORANGE)

add_text_block(slide,
    "Token-based billing with full transparency via ACCOUNT_USAGE views\n\n"
    "Two billing surfaces (tracked separately):\n"
    "  • CORTEX_CODE_CLI — usage from the terminal / IDE CLI\n"
    "  • CORTEX_CODE_SNOWSIGHT — usage from the Snowsight web interface\n\n"
    "What counts as CoCo credit usage: AI model inference (prompt + response tokens)\n"
    "  — Planning, code generation, data discovery conversations\n\n"
    "Billed separately (standard Snowflake rates):\n"
    "  • Warehouse compute for SQL execution (SELECT, CREATE, etc.)\n"
    "  • Storage for tables, stages, and dynamic table refreshes\n"
    "  • Cortex Agent / Analyst usage (separate AI_SERVICES metering)",
    0.6, 3.4, 12.0, Pt(13), WHITE)

add_progress_bar(slide, 8)

# ============================================================
# SLIDE 10: Cost Controls & Governance
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)
add_title(slide, "Cost Controls & Governance")

# 3-column: Limits → Monitoring → Best Practices
add_flow_box(slide, "Daily Limits", "Per-user credit caps\n(account or user level)", 0.6, 1.5, 3.5, 1.5, RED_ACCENT)
add_arrow(slide, 4.1, 2.0, 0.8)
add_flow_box(slide, "Usage Monitoring", "ACCOUNT_USAGE views\nwith full session detail", 4.9, 1.5, 3.5, 1.5, BLUE)
add_arrow(slide, 8.4, 2.0, 0.8)
add_flow_box(slide, "Governance", "RBAC controls + alerts\nfor proactive management", 9.2, 1.5, 3.5, 1.5, GREEN)

add_text_block(slide,
    "Daily Credit Limits (rolling 24-hour window):\n"
    "  • CORTEX_CODE_CLI_DAILY_EST_CREDIT_LIMIT_PER_USER\n"
    "  • CORTEX_CODE_SNOWSIGHT_DAILY_EST_CREDIT_LIMIT_PER_USER\n"
    "  • Value: -1 = unlimited, 0 = blocked, positive number = daily cap\n\n"
    "Usage Monitoring Views (SNOWFLAKE.ACCOUNT_USAGE):\n"
    "  • CORTEX_CODE_CLI_USAGE_HISTORY — per-session token & credit detail\n"
    "  • METERING_DAILY_HISTORY (service_type = 'CORTEX_CODE_CLI') — daily totals\n\n"
    "Best Practices:\n"
    "  • Set default daily limits (e.g., 20 credits/day), grant higher for power users\n"
    "  • Monitor weekly trends, combine with RBAC to control which roles can use CoCo",
    0.6, 3.4, 12.0, Pt(13), WHITE)

add_progress_bar(slide, 9)

# ============================================================
# SLIDE 11: Summary & Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_dark_bg(slide)
add_title(slide, "Summary & Next Steps")

# What we built — 3 boxes
add_flow_box(slide, "Built", "4-source AP pipeline\nDynamic Table", 0.6, 1.5, 3.5, 1.3, ORANGE)
add_flow_box(slide, "Created", "Reusable skill +\nSemantic View", 4.9, 1.5, 3.5, 1.3, PURPLE)
add_flow_box(slide, "Deployed", "Cortex Agent +\nStreamlit App", 9.2, 1.5, 3.5, 1.3, GREEN)

# Key takeaways
txBox = slide.shapes.add_textbox(Inches(0.6), Inches(3.1), Inches(6.0), Inches(0.4))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Key Takeaways"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = BLUE

add_text_block(slide,
    "Cortex Code accelerates data engineering with AI-assisted development\n"
    "Plan Mode + validation ensures safety and correctness at every step\n"
    "Skills make workflows repeatable and consistent across the team\n"
    "Semantic Views + Agents democratize data for business users",
    0.6, 3.6, 6.0, Pt(14), WHITE)

# Next steps
txBox2 = slide.shapes.add_textbox(Inches(7.0), Inches(3.1), Inches(5.5), Inches(0.4))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Next Steps"
p2.font.size = Pt(16)
p2.font.bold = True
p2.font.color.rgb = BLUE

add_text_block(slide,
    "Resolve payment terms normalization (OQ-1)\n"
    "Deploy Gold layer DT (TARGET_LAG = DOWNSTREAM)\n"
    "Evaluate Agent with 15-question golden set\n"
    "Roll out CoCo CLI to broader data team",
    7.0, 3.6, 5.5, Pt(14), WHITE)

add_progress_bar(slide, 10)

# ============================================================
# Save
# ============================================================
output_path = os.path.join(SCRIPT_DIR, "Cortex_Code_Demo_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
